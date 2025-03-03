import logging
from typing import List, Dict, Any, Optional
import os
import tempfile
import uuid
import asyncio
from concurrent.futures import ThreadPoolExecutor
import traceback
import time
import re

import chromadb
from chromadb.config import Settings
from django.conf import settings
from ollama import Client

from features.knowledge.repositories.knowledge_repository import KnowledgeRepository
from api.utils.exceptions import NotFoundException


class KnowledgeService:
    def __init__(self):
        self.repository = KnowledgeRepository()
        self.logger = logging.getLogger(__name__)
        self.ollama_client = Client(host=settings.OLLAMA_ENDPOINT)

        # Initialize ChromaDB with new client configuration
        self.chroma_client = chromadb.PersistentClient(
            path=settings.CHROMA_PERSIST_DIR,
            settings=Settings(anonymized_telemetry=False, allow_reset=False, is_persistent=True),
        )

        # Create or get collection with updated configuration
        self.collection = self.chroma_client.get_or_create_collection(
            name="knowledge_embeddings",
            metadata={"hnsw:space": "cosine"},
            embedding_function=None,  # We're using Ollama for embeddings
        )
        
        # Create executor for background processing
        self.executor = ThreadPoolExecutor(max_workers=2)
        
        # Initialize caches
        self._init_caches()

    def _init_caches(self):
        """Initialize caches for better performance"""
        # Cache for embeddings
        self._embedding_cache = {}
        self._embedding_cache_max_size = getattr(settings, 'EMBEDDING_CACHE_MAX_SIZE', 1000)
        
        # Cache for search results
        self._search_results_cache = {}
        self._search_results_cache_max_size = getattr(settings, 'SEARCH_RESULTS_CACHE_MAX_SIZE', 100)
        self._search_results_cache_ttl = getattr(settings, 'SEARCH_RESULTS_CACHE_TTL', 300)  # 5 minutes
        
        # Cache for chunks
        self._chunks_cache = {}
        self._chunks_cache_max_size = getattr(settings, 'CHUNKS_CACHE_MAX_SIZE', 50)
        
        # Track cache stats
        self._cache_hits = {'embedding': 0, 'search': 0, 'chunks': 0}
        self._cache_misses = {'embedding': 0, 'search': 0, 'chunks': 0}

    def create_knowledge_with_file(self, data: dict, user, file) -> Any:
        """
        Create a new knowledge document from a file.
        This method handles the initial creation with processing status,
        then processes the file in the background.
        
        Args:
            data: The knowledge document data
            user: The user creating the document
            file: The uploaded file object
            
        Returns:
            The created knowledge document
        """
        try:
            # First create the knowledge document with processing status
            data["user"] = user
            data["status"] = "processing"  # Explicitly set status to processing
            knowledge = self.repository.create(data)
            
            # Use a unique task ID to prevent duplicate processing
            task_id = f"process_file_{knowledge.id}"
            
            # Check if this task is already running
            if not hasattr(self, '_processing_tasks'):
                self._processing_tasks = set()
            
            if task_id in self._processing_tasks:
                self.logger.warning(f"File processing for knowledge {knowledge.id} is already in progress")
            else:
                # Mark this task as running
                self._processing_tasks.add(task_id)
                
                # Create a file-like object with the file content to avoid closed file issues
                file_info = {
                    'name': file.name,
                    'content_type': file.content_type,
                    'size': file.size
                }
                
                # Read the file content into memory first
                try:
                    if hasattr(file, 'read'):
                        file_content = file.read()
                        file_info['content'] = file_content
                    else:
                        # If file is already closed, try to open it from the path
                        try:
                            with open(file.temporary_file_path(), 'rb') as f:
                                file_content = f.read()
                                file_info['content'] = file_content
                        except (AttributeError, FileNotFoundError):
                            # If we can't get the file content, raise an error
                            raise ValueError(f"Could not read file content for {file.name}")
                except Exception as e:
                    self.logger.error(f"Error reading file content: {str(e)}")
                    raise
                
                # Process the file in the background with the file info instead of the file object
                future = self.executor.submit(self._process_file_with_content, knowledge, file_info)
                
                # Add a callback to remove the task ID when done
                def _task_done_callback(future):
                    try:
                        self._processing_tasks.remove(task_id)
                    except KeyError:
                        pass
                
                future.add_done_callback(_task_done_callback)
            
            return knowledge
        except Exception as e:
            self.logger.error(f"Error creating knowledge with file: {str(e)}")
            raise
    
    def _process_file_with_content(self, knowledge, file_info):
        """
        Process a file using its content that was already read into memory.
        
        Args:
            knowledge: The knowledge document to update
            file_info: Dictionary containing file information and content
        """
        try:
            self.logger.info(f"Processing file {file_info['name']} for knowledge {knowledge.id}")
            
            # Create a temporary file to store the file content
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                # Write the content to the temporary file
                temp_file.write(file_info['content'])
                temp_file_path = temp_file.name
            
            try:
                # Process the file based on its type
                content = ""
                chunks = []
                metadata = {"source": file_info['name']}
                
                file_type = file_info['content_type'].lower()
                file_name = file_info['name'].lower()
                
                # PDF processing
                if "pdf" in file_type or file_name.endswith('.pdf'):
                    content, chunks, metadata = self._process_pdf(temp_file_path, knowledge.name)
                
                # Word document processing
                elif "word" in file_type or file_name.endswith(('.docx', '.doc')):
                    content, chunks, metadata = self._process_docx(temp_file_path, knowledge.name)
                
                # PowerPoint processing
                elif "presentation" in file_type or file_name.endswith(('.pptx', '.ppt')):
                    content, chunks, metadata = self._process_pptx(temp_file_path, knowledge.name)
                
                # Excel processing
                elif "excel" in file_type or "spreadsheet" in file_type or file_name.endswith(('.xlsx', '.xls')):
                    content, chunks, metadata = self._process_excel(temp_file_path, knowledge.name)
                
                # Markdown processing
                elif "markdown" in file_type or file_name.endswith('.md'):
                    content, chunks, metadata = self._process_markdown(temp_file_path, knowledge.name)
                
                # HTML processing
                elif "html" in file_type or file_name.endswith(('.html', '.htm')):
                    content, chunks, metadata = self._process_html(temp_file_path, knowledge.name)
                
                # JSON processing
                elif "json" in file_type or file_name.endswith('.json'):
                    content, chunks, metadata = self._process_json(temp_file_path, knowledge.name)
                
                # CSV processing
                elif "csv" in file_type or file_name.endswith('.csv'):
                    content, chunks, metadata = self._process_csv(temp_file_path)
                
                # Plain text processing
                elif "text" in file_type or file_name.endswith('.txt'):
                    with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    text_chunks = self._chunk_text(content)
                    # Convert text chunks to the expected dictionary format
                    chunks = []
                    for i, chunk_text in enumerate(text_chunks):
                        chunk_id = f"{file_info['name']}_c{i}"
                        chunks.append({
                            "id": chunk_id,
                            "content": chunk_text,
                            "metadata": {
                                "source": file_info['name'],
                                "chunk": i,
                                "citation": file_info['name']
                            }
                        })
                
                # Default processing for other file types
                else:
                    with open(temp_file_path, 'rb') as f:
                        binary_content = f.read()
                    try:
                        content = binary_content.decode('utf-8', errors='ignore')
                    except:
                        content = f"Binary file content could not be processed: {file_info['name']}"
                    text_chunks = self._chunk_text(content)
                    # Convert text chunks to the expected dictionary format
                    chunks = []
                    for i, chunk_text in enumerate(text_chunks):
                        chunk_id = f"{file_info['name']}_c{i}"
                        chunks.append({
                            "id": chunk_id,
                            "content": chunk_text,
                            "metadata": {
                                "source": file_info['name'],
                                "chunk": i,
                                "citation": file_info['name']
                            }
                        })
                
                # Update the knowledge document with the processed content
                self.repository.update(knowledge.id, {
                    "content": content,
                    "status": "ready",
                })
                
                # Store chunks in ChromaDB with metadata for citation
                self._store_chunks_in_chroma(knowledge.id, chunks, metadata, knowledge.user_id)
                
                self.logger.info(f"Successfully processed file for knowledge {knowledge.id}")
            finally:
                # Clean up the temporary file
                os.unlink(temp_file_path)
                
        except Exception as e:
            self.logger.error(f"Error processing file: {str(e)}")
            traceback.print_exc()
            # Update knowledge document with error status
            self.repository.update(knowledge.id, {
                "status": "error",
                "error_message": str(e),
            })
    
    def _process_pdf(self, file_path, document_name):
        """
        Process a PDF file and extract its content with page information.
        
        Args:
            file_path: Path to the PDF file
            document_name: Name of the document
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            # Import PyPDF2 here to avoid dependency issues
            import PyPDF2
            
            full_content = ""
            chunks = []
            metadata = {"source": document_name, "type": "pdf"}
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata["total_pages"] = len(pdf_reader.pages)
                
                for i, page in enumerate(pdf_reader.pages):
                    page_num = i + 1
                    page_text = page.extract_text()
                    
                    if page_text:
                        # Add page number to the content
                        page_content = f"Page {page_num}:\n{page_text}\n\n"
                        full_content += page_content
                        
                        # Create chunks from the page with citation metadata
                        page_chunks = self._chunk_text(page_text)
                        for j, chunk in enumerate(page_chunks):
                            chunk_id = f"{document_name}_p{page_num}_c{j}"
                            chunks.append({
                                "id": chunk_id,
                                "content": chunk,
                                "metadata": {
                                    "source": document_name,
                                    "page": page_num,
                                    "chunk": j,
                                    "citation": f"{document_name}, Page {page_num}"
                                }
                            })
            
            return full_content, chunks, metadata
        except ImportError:
            self.logger.error("PyPDF2 is not installed. Please install it to process PDF files.")
            raise Exception("PDF processing library not available")
        except Exception as e:
            self.logger.error(f"Error processing PDF: {str(e)}")
            raise
    
    def _process_csv(self, file_path):
        """
        Process a CSV file and extract its content.
        
        Args:
            file_path: Path to the CSV file
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            import csv
            
            full_content = ""
            chunks = []
            metadata = {"source": os.path.basename(file_path), "type": "csv"}
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                csv_reader = csv.reader(file)
                rows = list(csv_reader)
                
                # Get headers if available
                headers = rows[0] if rows else []
                metadata["headers"] = headers
                
                # Convert CSV to text format
                for i, row in enumerate(rows):
                    row_text = ", ".join(row)
                    full_content += f"Row {i+1}: {row_text}\n"
                    
                    # Create chunk for each row with citation metadata
                    chunk_id = f"{metadata['source']}_r{i+1}"
                    chunks.append({
                        "id": chunk_id,
                        "content": row_text,
                        "metadata": {
                            "source": metadata['source'],
                            "row": i+1,
                            "citation": f"{metadata['source']}, Row {i+1}"
                        }
                    })
            
            return full_content, chunks, metadata
        except Exception as e:
            self.logger.error(f"Error processing CSV: {str(e)}")
            raise
    
    def _process_docx(self, file_path, document_name):
        """
        Process a Word document (DOCX/DOC) and extract its content.
        
        Args:
            file_path: Path to the Word document
            document_name: Name of the document
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            # Import docx library here to avoid dependency issues
            import docx
            
            full_content = ""
            chunks = []
            metadata = {"source": document_name, "type": "docx"}
            
            # Open the document
            doc = docx.Document(file_path)
            
            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():  # Skip empty paragraphs
                    paragraphs.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                    if row_text.strip():
                        paragraphs.append(row_text)
            
            # Combine all paragraphs
            full_content = "\n\n".join(paragraphs)
            
            # Create chunks with citation metadata
            text_chunks = self._chunk_text(full_content)
            for i, chunk in enumerate(text_chunks):
                chunk_id = f"{document_name}_c{i}"
                chunks.append({
                    "id": chunk_id,
                    "content": chunk,
                    "metadata": {
                        "source": document_name,
                        "chunk": i,
                        "citation": document_name
                    }
                })
            
            return full_content, chunks, metadata
        except ImportError:
            self.logger.error("python-docx is not installed. Please install it to process DOCX files.")
            raise Exception("DOCX processing library not available")
        except Exception as e:
            self.logger.error(f"Error processing DOCX: {str(e)}")
            raise
    
    def _process_excel(self, file_path, document_name):
        """
        Process an Excel file (XLSX/XLS) and extract its content.
        
        Args:
            file_path: Path to the Excel file
            document_name: Name of the document
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            # Import pandas and openpyxl here to avoid dependency issues
            import pandas as pd
            
            full_content = ""
            chunks = []
            metadata = {"source": document_name, "type": "excel"}
            
            # Read all sheets from the Excel file
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            metadata["sheets"] = sheet_names
            
            # Process each sheet
            for sheet_name in sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert sheet to text
                sheet_content = f"Sheet: {sheet_name}\n"
                sheet_content += df.to_string(index=False) + "\n\n"
                full_content += sheet_content
                
                # Create chunks for each sheet with citation metadata
                sheet_chunks = self._chunk_text(sheet_content)
                for i, chunk in enumerate(sheet_chunks):
                    chunk_id = f"{document_name}_{sheet_name}_c{i}"
                    chunks.append({
                        "id": chunk_id,
                        "content": chunk,
                        "metadata": {
                            "source": document_name,
                            "sheet": sheet_name,
                            "chunk": i,
                            "citation": f"{document_name}, Sheet: {sheet_name}"
                        }
                    })
            
            return full_content, chunks, metadata
        except ImportError:
            self.logger.error("pandas or openpyxl is not installed. Please install them to process Excel files.")
            raise Exception("Excel processing libraries not available")
        except Exception as e:
            self.logger.error(f"Error processing Excel: {str(e)}")
            raise
    
    def _process_markdown(self, file_path, document_name):
        """
        Process a Markdown file and extract its content.
        
        Args:
            file_path: Path to the Markdown file
            document_name: Name of the document
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            full_content = ""
            chunks = []
            metadata = {"source": document_name, "type": "markdown"}
            
            # Read the markdown file
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                full_content = f.read()
            
            # Create chunks with citation metadata
            text_chunks = self._chunk_text(full_content)
            for i, chunk in enumerate(text_chunks):
                chunk_id = f"{document_name}_c{i}"
                chunks.append({
                    "id": chunk_id,
                    "content": chunk,
                    "metadata": {
                        "source": document_name,
                        "chunk": i,
                        "citation": document_name
                    }
                })
            
            return full_content, chunks, metadata
        except Exception as e:
            self.logger.error(f"Error processing Markdown: {str(e)}")
            raise
    
    def _process_html(self, file_path, document_name):
        """
        Process an HTML file and extract its content.
        
        Args:
            file_path: Path to the HTML file
            document_name: Name of the document
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            # Import BeautifulSoup here to avoid dependency issues
            from bs4 import BeautifulSoup
            
            full_content = ""
            chunks = []
            metadata = {"source": document_name, "type": "html"}
            
            # Read the HTML file
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            
            # Extract text
            text = soup.get_text()
            
            # Clean up text (remove extra whitespace)
            lines = (line.strip() for line in text.splitlines())
            chunks_of_lines = (phrase.strip() for line in lines for phrase in line.split("  "))
            full_content = '\n'.join(chunk for chunk in chunks_of_lines if chunk)
            
            # Extract title if available
            title_tag = soup.find('title')
            if title_tag:
                metadata["title"] = title_tag.string
            
            # Create chunks with citation metadata
            text_chunks = self._chunk_text(full_content)
            for i, chunk in enumerate(text_chunks):
                chunk_id = f"{document_name}_c{i}"
                chunks.append({
                    "id": chunk_id,
                    "content": chunk,
                    "metadata": {
                        "source": document_name,
                        "chunk": i,
                        "citation": metadata.get("title", document_name)
                    }
                })
            
            return full_content, chunks, metadata
        except ImportError:
            self.logger.error("BeautifulSoup is not installed. Please install it to process HTML files.")
            raise Exception("HTML processing library not available")
        except Exception as e:
            self.logger.error(f"Error processing HTML: {str(e)}")
            raise
            
    def _process_json(self, file_path, document_name):
        """
        Process a JSON file and extract its content.
        
        Args:
            file_path: Path to the JSON file
            document_name: Name of the document
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            import json
            
            full_content = ""
            chunks = []
            metadata = {"source": document_name, "type": "json"}
            
            # Read the JSON file
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                json_data = json.load(f)
            
            # Convert JSON to formatted string
            full_content = json.dumps(json_data, indent=2)
            
            # For large JSON files, process by sections
            if len(full_content) > 10000:  # If JSON is large
                # Process by top-level keys
                for key, value in json_data.items():
                    section = f"{key}:\n{json.dumps(value, indent=2)}"
                    section_chunks = self._chunk_text(section)
                    
                    for i, chunk in enumerate(section_chunks):
                        chunk_id = f"{document_name}_{key}_c{i}"
                        chunks.append({
                            "id": chunk_id,
                            "content": chunk,
                            "metadata": {
                                "source": document_name,
                                "section": key,
                                "chunk": i,
                                "citation": f"{document_name}, Section: {key}"
                            }
                        })
            else:
                # For smaller JSON, just chunk the whole thing
                text_chunks = self._chunk_text(full_content)
                for i, chunk in enumerate(text_chunks):
                    chunk_id = f"{document_name}_c{i}"
                    chunks.append({
                        "id": chunk_id,
                        "content": chunk,
                        "metadata": {
                            "source": document_name,
                            "chunk": i,
                            "citation": document_name
                        }
                    })
            
            return full_content, chunks, metadata
        except Exception as e:
            self.logger.error(f"Error processing JSON: {str(e)}")
            raise
    
    def _chunk_text(self, text, chunk_size=1000, overlap=100):
        """
        Split text into overlapping chunks based on semantic boundaries.
        This improved chunking strategy respects paragraph and sentence boundaries
        for more coherent chunks.
        
        Args:
            text: The text to chunk
            chunk_size: Maximum size of each chunk
            overlap: Number of characters to overlap between chunks
            
        Returns:
            List of text chunks
        """
        if not text:
            return []
            
        # Use semantic chunking strategy based on settings
        chunking_strategy = getattr(settings, 'CHUNKING_STRATEGY', 'semantic')
        
        if chunking_strategy == 'simple':
            return self._chunk_text_simple(text, chunk_size, overlap)
        else:
            return self._chunk_text_semantic(text, chunk_size, overlap)
            
    def _chunk_text_simple(self, text, chunk_size=1000, overlap=100):
        """
        Simple chunking strategy that splits text by character count.
        
        Args:
            text: The text to chunk
            chunk_size: Maximum size of each chunk
            overlap: Number of characters to overlap between chunks
            
        Returns:
            List of text chunks
        """
        chunks = []
        
        # Simple chunking with overlap
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            
            # If this is not the first chunk and we're not at the end, 
            # try to find a good breaking point (space or newline)
            if start > 0 and end < len(text):
                # Look for a space or newline to break at
                break_point = text.rfind(' ', end - 50, end)
                if break_point == -1:
                    break_point = text.rfind('\n', end - 50, end)
                
                if break_point != -1:
                    end = break_point + 1  # Include the space or newline
            
            # Add the chunk
            chunks.append(text[start:end])
            
            # Move start position for next chunk, considering overlap
            start = end - overlap if end < len(text) else len(text)
        
        return chunks
            
    def _chunk_text_semantic(self, text, chunk_size=1000, overlap=100):
        """
        Advanced chunking strategy that respects semantic boundaries like
        paragraphs, sentences, and sections.
        
        Args:
            text: The text to chunk
            chunk_size: Maximum size of each chunk
            overlap: Number of characters to overlap between chunks
            
        Returns:
            List of text chunks
        """
        chunks = []
        
        # First split by paragraphs (double newlines)
        paragraphs = re.split(r'\n\s*\n', text)
        
        # Then split paragraphs into sentences
        sentence_pattern = r'(?<=[.!?])\s+'
        
        current_chunk = ""
        last_sentences = []  # Keep track of recent sentences for overlap
        
        for paragraph in paragraphs:
            # Skip empty paragraphs
            if not paragraph.strip():
                continue
                
            # Split paragraph into sentences
            sentences = re.split(sentence_pattern, paragraph)
            
            for sentence in sentences:
                # Skip empty sentences
                if not sentence.strip():
                    continue
                    
                # If adding this sentence would exceed chunk size, save current chunk and start a new one
                if len(current_chunk) + len(sentence) + 1 > chunk_size and current_chunk:
                    chunks.append(current_chunk)
                    
                    # Start new chunk with overlap from the end of the previous chunk
                    # Use the last few sentences for overlap
                    overlap_text = ""
                    overlap_size = 0
                    
                    # Add sentences from the end until we reach desired overlap
                    for prev_sentence in reversed(last_sentences):
                        if overlap_size + len(prev_sentence) > overlap:
                            break
                        overlap_text = prev_sentence + " " + overlap_text
                        overlap_size += len(prev_sentence) + 1
                    
                    current_chunk = overlap_text + sentence
                    # Reset last_sentences but keep the current sentence
                    last_sentences = [sentence]
                else:
                    # Add sentence to current chunk with a space if needed
                    if current_chunk and not current_chunk.endswith(" "):
                        current_chunk += " "
                    current_chunk += sentence
                    
                    # Keep track of recent sentences for overlap
                    last_sentences.append(sentence)
                    # Only keep enough sentences for overlap
                    while sum(len(s) + 1 for s in last_sentences) > overlap * 2:
                        last_sentences.pop(0)
            
            # Add paragraph break if not at the end of a chunk
            if current_chunk and not current_chunk.endswith("\n"):
                current_chunk += "\n\n"
        
        # Add the last chunk if it's not empty
        if current_chunk:
            chunks.append(current_chunk)
            
        # Handle case where chunks are still too large (fallback)
        final_chunks = []
        for chunk in chunks:
            if len(chunk) > chunk_size * 1.5:  # If chunk is significantly larger than target
                # Fall back to simple chunking for this large chunk
                sub_chunks = self._chunk_text_simple(chunk, chunk_size, overlap)
                final_chunks.extend(sub_chunks)
            else:
                final_chunks.append(chunk)
                
        return final_chunks
    
    def _store_chunks_in_chroma(self, knowledge_id, chunks, metadata, user_id):
        """
        Store document chunks in ChromaDB with metadata for retrieval and citation.
        
        Args:
            knowledge_id: ID of the knowledge document
            chunks: List of chunks with content and metadata
            metadata: General metadata for the document
            user_id: ID of the user who owns the document
        """
        try:
            if not chunks:
                self.logger.warning(f"No chunks to store for knowledge {knowledge_id}")
                return
                
            # Log chunk information for debugging
            self.logger.debug(f"Storing {len(chunks)} chunks for knowledge {knowledge_id}")
            if chunks and isinstance(chunks[0], str):
                self.logger.warning(f"Chunks are strings, not dictionaries. Converting to proper format.")
                formatted_chunks = []
                for i, chunk_text in enumerate(chunks):
                    chunk_id = f"auto_{i}"
                    formatted_chunks.append({
                        "id": chunk_id,
                        "content": chunk_text,
                        "metadata": {
                            "source": metadata.get("source", "unknown"),
                            "chunk": i,
                            "citation": metadata.get("source", "unknown")
                        }
                    })
                chunks = formatted_chunks
                self.logger.debug(f"Converted {len(chunks)} string chunks to dictionary format")
            
            # Generate embeddings for all chunks
            embeddings = []
            ids = []
            documents = []
            metadatas = []
            
            for chunk in chunks:
                # Create a unique ID for this chunk
                try:
                    # Handle case where chunk might be a string instead of a dictionary
                    if isinstance(chunk, str):
                        self.logger.warning(f"Found string chunk instead of dictionary. Converting to proper format.")
                        chunk_text = chunk
                        chunk = {
                            "id": f"auto_{uuid.uuid4()}",
                            "content": chunk_text,
                            "metadata": {
                                "source": metadata.get("source", "unknown"),
                                "citation": metadata.get("source", "unknown")
                            }
                        }
                    
                    chunk_id = f"{knowledge_id}_{chunk['id']}" if 'id' in chunk else f"{knowledge_id}_{uuid.uuid4()}"
                    
                    # Log chunk information for debugging
                    self.logger.debug(f"Processing chunk: {chunk_id}")
                    self.logger.debug(f"Chunk type: {type(chunk)}")
                    self.logger.debug(f"Chunk keys: {chunk.keys() if isinstance(chunk, dict) else 'Not a dictionary'}")
                    
                    chunk_embedding = self._generate_embedding(chunk["content"])
                    
                    # Combine chunk metadata with document metadata
                    chunk_metadata = chunk.get("metadata", {})
                    combined_metadata = {
                        "user_id": str(user_id),
                        "knowledge_id": str(knowledge_id),
                        "source": metadata.get("source", ""),
                        **chunk_metadata
                    }
                    
                    embeddings.append(chunk_embedding)
                    ids.append(chunk_id)
                    documents.append(chunk["content"])
                    metadatas.append(combined_metadata)
                except Exception as e:
                    self.logger.error(f"Error generating embedding for chunk: {str(e)}")
                    self.logger.error(f"Chunk data: {chunk}")
                    traceback.print_exc()
                    continue
            
            if not embeddings:
                self.logger.warning(f"No valid embeddings generated for knowledge {knowledge_id}")
                return
                
            # Store in ChromaDB in batches to avoid memory issues
            batch_size = 50
            for i in range(0, len(embeddings), batch_size):
                end_idx = min(i + batch_size, len(embeddings))
                self.collection.add(
                    ids=ids[i:end_idx],
                    embeddings=embeddings[i:end_idx],
                    documents=documents[i:end_idx],
                    metadatas=metadatas[i:end_idx],
                )
            
            self.logger.info(f"Stored {len(chunks)} chunks for knowledge {knowledge_id} in ChromaDB")
        except Exception as e:
            self.logger.error(f"Error storing chunks in ChromaDB: {str(e)}")
            traceback.print_exc()

    def create_knowledge(self, data: dict, user):
        """Create a new knowledge document"""
        try:
            # First create in Django for metadata
            data["user"] = user
            knowledge = self.repository.create(data)

            # Generate embedding and store in ChromaDB
            embedding = self._generate_embedding(data["content"])
            
            # Store in ChromaDB
            self.collection.add(
                ids=[str(knowledge.id)],
                embeddings=[embedding],
                documents=[data["content"]],
                metadatas=[
                    {
                        "user_id": str(user.id),
                        "knowledge_id": str(knowledge.id),
                        "name": data["name"],
                        "identifier": data["identifier"],
                        "citation": data["name"]
                    }
                ],
            )
            return knowledge
        except Exception as e:
            self.logger.error(f"Error creating knowledge: {str(e)}")
            raise

    def get_knowledge(self, knowledge_id, user_id=None):
        """
        Get knowledge by ID.
        
        Args:
            knowledge_id: The UUID of the knowledge to get
            user_id: Optional user ID to filter by
            
        Returns:
            Knowledge object if found, None otherwise
        """
        try:
            print(f"DEBUG: KnowledgeService.get_knowledge: Getting knowledge with ID {knowledge_id}, user_id: {user_id}")
            # Pass the knowledge_id directly to the repository
            knowledge = self.repository.get_by_id(knowledge_id, user_id)
            if knowledge:
                print(f"DEBUG: KnowledgeService.get_knowledge: Found knowledge: {knowledge.name}")
            else:
                print(f"DEBUG: KnowledgeService.get_knowledge: Knowledge not found for ID: {knowledge_id}")
            return knowledge
        except Exception as e:
            print(f"DEBUG: KnowledgeService.get_knowledge: Error getting knowledge: {str(e)}")
            self.logger.error(f"Error getting knowledge: {str(e)}")
            return None

    def get_by_identifier(self, identifier: str, user_id: int):
        """Get knowledge by identifier"""
        knowledge = self.repository.get_by_identifier(identifier, user_id)
        if not knowledge:
            raise NotFoundException("Knowledge document not found")
        return knowledge

    def list_knowledge(self, user_id: int):
        """List knowledge documents"""
        return self.repository.get_user_knowledge(user_id)

    def update_knowledge(self, knowledge_id: int, data: dict, user_id: int):
        """Update knowledge document and its embedding"""
        try:
            # Verify ownership and get knowledge
            knowledge = self.get_knowledge(knowledge_id, user_id)

            # Update Django model
            updated = self.repository.update(knowledge.id, data)

            # Update ChromaDB if content changed
            if "content" in data:
                embedding = self._generate_embedding(data["content"])
                
                # First delete existing entries for this knowledge document
                self.collection.delete(
                    where={"knowledge_id": str(knowledge_id)}
                )
                
                # Then add the updated content
                self.collection.add(
                    ids=[str(knowledge_id)],
                    embeddings=[embedding],
                    documents=[data["content"]],
                    metadatas=[
                        {
                            "user_id": str(user_id),
                            "knowledge_id": str(knowledge_id),
                            "name": data.get("name", knowledge.name),
                            "identifier": data.get("identifier", knowledge.identifier),
                            "citation": data.get("name", knowledge.name)
                        }
                    ],
                )

            return updated
        except Exception as e:
            self.logger.error(f"Error updating knowledge: {str(e)}")
            raise

    def delete_knowledge(self, knowledge_id: int, user_id: int):
        """Delete knowledge document from both Django and ChromaDB"""
        try:
            # Verify ownership and get knowledge
            knowledge = self.get_knowledge(knowledge_id, user_id)

            # Delete from both storages
            self.repository.delete(knowledge.id)
            
            # Delete all chunks associated with this knowledge document
            self.collection.delete(
                where={"knowledge_id": str(knowledge_id)}
            )
            
            return True
        except Exception as e:
            self.logger.error(f"Error deleting knowledge: {str(e)}")
            raise

    def _generate_embedding(self, text: str) -> List[float]:
        """
        Generate an embedding for the given text using either Ollama or sentence-transformers.
        Uses caching to avoid regenerating embeddings for the same text.
        
        Args:
            text: The text to generate an embedding for
            
        Returns:
            List of floats representing the embedding
        """
        if not text or len(text.strip()) == 0:
            self.logger.warning("Empty text provided for embedding generation")
            return []
            
        try:
            # Check cache first
            # Use a hash of the text as the cache key to avoid storing large strings
            cache_key = hash(text)
            if cache_key in self._embedding_cache:
                self._cache_hits['embedding'] += 1
                return self._embedding_cache[cache_key]
            
            self._cache_misses['embedding'] += 1
            
            # Truncate text if it's too long to avoid token limits
            # Most embedding models have limits around 8192 tokens
            max_chars = 8000  # Approximate character limit
            if len(text) > max_chars:
                self.logger.info(f"Truncating text from {len(text)} to {max_chars} characters for embedding")
                text = text[:max_chars]
            
            # Check if we should use sentence-transformers (if available)
            use_sentence_transformers = getattr(settings, 'USE_SENTENCE_TRANSFORMERS', False)
            
            if use_sentence_transformers:
                try:
                    # Import here to avoid dependency issues if not installed
                    from sentence_transformers import SentenceTransformer
                    
                    # Use a singleton pattern for the model to avoid reloading
                    if not hasattr(self, '_sentence_transformer_model'):
                        model_name = getattr(settings, 'SENTENCE_TRANSFORMER_MODEL', 'all-MiniLM-L6-v2')
                        self.logger.info(f"Loading sentence-transformer model: {model_name}")
                        self._sentence_transformer_model = SentenceTransformer(model_name)
                    
                    # Generate embedding
                    embedding = self._sentence_transformer_model.encode(text, convert_to_numpy=True).tolist()
                    self.logger.debug(f"Generated embedding with sentence-transformers, dimension: {len(embedding)}")
                    
                    # Cache the result
                    self._embedding_cache[cache_key] = embedding
                    
                    # Manage cache size
                    if len(self._embedding_cache) > self._embedding_cache_max_size:
                        # Remove a random key to avoid complex LRU logic
                        self._embedding_cache.pop(next(iter(self._embedding_cache)))
                        
                    return embedding
                except ImportError:
                    self.logger.warning("sentence-transformers not available, falling back to Ollama")
                except Exception as e:
                    self.logger.error(f"Error using sentence-transformers: {str(e)}")
                    self.logger.warning("Falling back to Ollama for embeddings")
            
            # Use Ollama's embedding endpoint as fallback
            embedding_model = getattr(settings, 'EMBEDDING_MODEL', 'nomic-embed-text')
            response = self.ollama_client.embeddings(
                model=embedding_model,
                prompt=text,
            )
            
            if not response or "embedding" not in response:
                self.logger.error(f"Failed to generate embedding: {response}")
                return []
                
            embedding = response["embedding"]
            
            # Validate embedding
            if not embedding or len(embedding) == 0:
                self.logger.error("Empty embedding returned from Ollama")
                return []
                
            # Check if embedding is all zeros or very small values
            if all(abs(v) < 1e-6 for v in embedding):
                self.logger.warning("Embedding contains all zeros or very small values")
            
            # Cache the result
            self._embedding_cache[cache_key] = embedding
            
            # Manage cache size
            if len(self._embedding_cache) > self._embedding_cache_max_size:
                # Remove a random key to avoid complex LRU logic
                self._embedding_cache.pop(next(iter(self._embedding_cache)))
            
            return embedding
        except Exception as e:
            self.logger.error(f"Error generating embedding: {str(e)}")
            traceback.print_exc()
            return []

    def find_relevant_context(self, query: str, user_id: int, max_results: int = 3) -> List[dict]:
        """Find relevant knowledge documents using hybrid search (semantic + keyword)"""
        try:
            if not query or len(query.strip()) == 0:
                self.logger.warning("Empty query provided to find_relevant_context")
                return []
            
            # Check if the user has any knowledge documents
            user_knowledge = self.repository.get_user_knowledge(user_id)
            if not user_knowledge or len(user_knowledge) == 0:
                self.logger.info(f"User {user_id} has no knowledge documents")
                return []
            
            # Determine if we should use hybrid search
            use_hybrid_search = getattr(settings, 'USE_HYBRID_SEARCH', True)
            
            if use_hybrid_search:
                self.logger.info(f"Using hybrid search for query: {query[:50]}...")
                return self._hybrid_search(query, user_id, max_results)
            else:
                self.logger.info(f"Using semantic search for query: {query[:50]}...")
                return self._semantic_search(query, user_id, max_results)
        except Exception as e:
            self.logger.error(f"Error finding relevant context: {str(e)}")
            traceback.print_exc()
            return []
            
    def _semantic_search(self, query: str, user_id: int, max_results: int = 3) -> List[dict]:
        """Find relevant knowledge documents using semantic search"""
        try:
            self.logger.info(f"Generating embedding for query: {query[:50]}...")
            query_embedding = self._generate_embedding(query)
            
            # Check if we have a valid embedding
            if not query_embedding or all(v == 0 for v in query_embedding):
                self.logger.warning("Failed to generate valid embedding for query")
                return []
            
            # Query ChromaDB for relevant chunks
            self.logger.info(f"Querying ChromaDB for user {user_id} with query: {query[:50]}...")
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=max_results * 3,  # Get more results than needed to filter
                where={"user_id": str(user_id)},
                include=["documents", "metadatas", "distances"],
            )
            
            if not results or not results["documents"] or len(results["documents"][0]) == 0:
                self.logger.info(f"No relevant context found for query: {query[:50]}...")
                return []

            # Process and format results
            formatted_results = []
            for i, (doc, meta, dist) in enumerate(zip(
                results["documents"][0], results["metadatas"][0], results["distances"][0]
            )):
                # Generate a unique ID for the chunk
                chunk_id = f"chunk_{meta.get('knowledge_id', '')}_{i}"
                
                # Calculate similarity score (1 - distance)
                similarity = 1 - dist
                
                # Only include results with reasonable similarity
                if similarity > 0.5:  # Adjust threshold as needed
                    formatted_results.append({
                        "id": chunk_id,  # Include the chunk ID
                        "content": doc,
                        "metadata": {
                            **meta,
                            "citation": meta.get("citation", meta.get("source", "Unknown Source"))
                        },
                        "similarity": similarity,
                        "search_type": "semantic"
                    })
            
            # Sort by similarity and limit to max_results
            formatted_results.sort(key=lambda x: x["similarity"], reverse=True)
            result_count = len(formatted_results[:max_results])
            self.logger.info(f"Found {result_count} relevant documents with semantic search")
            
            return formatted_results[:max_results]
        except Exception as e:
            self.logger.error(f"Error in semantic search: {str(e)}")
            traceback.print_exc()
            return []
            
    def _keyword_search(self, query: str, user_id: int, max_results: int = 3) -> List[dict]:
        """Find relevant knowledge documents using keyword search"""
        try:
            # Get all knowledge documents for the user
            user_knowledge = self.repository.get_user_knowledge(user_id)
            if not user_knowledge:
                return []
                
            # Extract keywords from query (simple approach)
            # Remove common stop words and split by spaces
            stop_words = {'a', 'an', 'the', 'and', 'or', 'but', 'is', 'are', 'was', 'were', 
                         'in', 'on', 'at', 'to', 'for', 'with', 'by', 'about', 'like', 
                         'through', 'over', 'before', 'after', 'between', 'under'}
            
            # Convert query to lowercase and tokenize
            query_lower = query.lower()
            query_tokens = [token.strip() for token in query_lower.split() if token.strip() and token.strip().lower() not in stop_words]
            
            # Get chunks for all knowledge documents
            all_chunks = []
            for knowledge in user_knowledge:
                chunks = self.get_chunks_for_knowledge(knowledge.id)
                all_chunks.extend(chunks)
                
            # Score chunks based on keyword matches
            scored_chunks = []
            for chunk in all_chunks:
                content = chunk.get('content', '').lower()
                metadata = chunk.get('metadata', {})
                
                # Count keyword occurrences
                keyword_count = sum(content.count(keyword) for keyword in query_tokens)
                
                # Calculate a simple TF score
                score = keyword_count / (len(content.split()) + 1) if content else 0
                
                # Only include chunks with at least one keyword match
                if score > 0:
                    scored_chunks.append({
                        "id": chunk.get('id', f"chunk_{len(scored_chunks)}"),
                        "content": chunk.get('content', ''),
                        "metadata": metadata,
                        "similarity": score,
                        "search_type": "keyword"
                    })
            
            # Sort by score and limit results
            scored_chunks.sort(key=lambda x: x["similarity"], reverse=True)
            result_count = len(scored_chunks[:max_results])
            self.logger.info(f"Found {result_count} relevant documents with keyword search")
            
            return scored_chunks[:max_results]
        except Exception as e:
            self.logger.error(f"Error in keyword search: {str(e)}")
            traceback.print_exc()
            return []
            
    def _hybrid_search(self, query: str, user_id: int, max_results: int = 3) -> List[dict]:
        """
        Combine semantic and keyword search results for better retrieval.
        Uses caching to improve performance for repeated queries.
        """
        try:
            # Generate cache key
            cache_key = f"search_{hash(query)}_{user_id}_{max_results}"
            
            # Check if we have a cached result that's still valid
            if cache_key in self._search_results_cache:
                cache_entry = self._search_results_cache[cache_key]
                cache_time = cache_entry.get('time', 0)
                current_time = time.time()
                
                # Check if the cache entry is still valid
                if current_time - cache_time < self._search_results_cache_ttl:
                    self._cache_hits['search'] += 1
                    return cache_entry.get('results', [])
            
            self._cache_misses['search'] += 1
            
            # Get results from both search methods
            semantic_results = self._semantic_search(query, user_id, max_results)
            keyword_results = self._keyword_search(query, user_id, max_results)
            
            # Combine results, avoiding duplicates
            combined_results = []
            seen_content = set()
            
            # First add semantic results (they usually have higher quality)
            for result in semantic_results:
                content_hash = hash(result["content"])
                if content_hash not in seen_content:
                    seen_content.add(content_hash)
                    combined_results.append(result)
            
            # Then add keyword results that aren't duplicates
            for result in keyword_results:
                content_hash = hash(result["content"])
                if content_hash not in seen_content:
                    seen_content.add(content_hash)
                    combined_results.append(result)
            
            # Re-rank combined results
            # For simplicity, we'll normalize scores and use a weighted combination
            # Semantic search gets higher weight (0.7) than keyword search (0.3)
            for result in combined_results:
                if result["search_type"] == "semantic":
                    result["combined_score"] = result["similarity"] * 0.7
                else:
                    result["combined_score"] = result["similarity"] * 0.3
            
            # Sort by combined score
            combined_results.sort(key=lambda x: x["combined_score"], reverse=True)
            
            # Limit to max_results
            final_results = combined_results[:max_results]
            self.logger.info(f"Hybrid search returned {len(final_results)} results")
            
            # Cache the results
            self._search_results_cache[cache_key] = {
                'results': final_results,
                'time': time.time()
            }
            
            # Manage cache size
            if len(self._search_results_cache) > self._search_results_cache_max_size:
                # Find and remove the oldest entry
                oldest_key = min(
                    self._search_results_cache.keys(),
                    key=lambda k: self._search_results_cache[k].get('time', 0)
                )
                self._search_results_cache.pop(oldest_key)
            
            return final_results
        except Exception as e:
            self.logger.error(f"Error in hybrid search: {str(e)}")
            traceback.print_exc()
            # Fall back to semantic search if hybrid fails
            return self._semantic_search(query, user_id, max_results)

    def bulk_create_knowledge(self, data_list: List[dict], user):
        """Bulk create knowledge documents"""
        try:
            # First create all documents in Django
            knowledge_docs = []
            for data in data_list:
                data["user"] = user
                knowledge = self.repository.create(data)
                knowledge_docs.append(knowledge)

            # Generate embeddings and store in ChromaDB
            for i, (knowledge, data) in enumerate(zip(knowledge_docs, data_list)):
                try:
                    embedding = self._generate_embedding(data["content"])
                    
                    self.collection.add(
                        ids=[str(knowledge.id)],
                        embeddings=[embedding],
                        documents=[data["content"]],
                        metadatas=[
                            {
                                "user_id": str(user.id),
                                "knowledge_id": str(knowledge.id),
                                "name": data["name"],
                                "identifier": data["identifier"],
                                "citation": data["name"]
                            }
                        ],
                    )
                except Exception as e:
                    self.logger.error(f"Error processing document {i} in bulk create: {str(e)}")
                    # Continue with other documents

            return knowledge_docs
        except Exception as e:
            self.logger.error(f"Error in bulk create knowledge: {str(e)}")
            raise

    async def reindex_all_knowledge(self):
        """Reindex all knowledge documents in ChromaDB"""
        try:
            # Clear the collection
            self.collection.delete(where={})

            # Get all knowledge documents
            all_knowledge = self.repository.list()

            # Batch process to avoid memory issues
            batch_size = 50
            for i in range(0, len(all_knowledge), batch_size):
                batch = all_knowledge[i : i + batch_size]
                
                for knowledge in batch:
                    try:
                        embedding = self._generate_embedding(knowledge.content)
                        
                        self.collection.add(
                            ids=[str(knowledge.id)],
                            embeddings=[embedding],
                            documents=[knowledge.content],
                            metadatas=[
                                {
                                    "user_id": str(knowledge.user_id),
                                    "knowledge_id": str(knowledge.id),
                                    "name": knowledge.name,
                                    "identifier": knowledge.identifier,
                                    "citation": knowledge.name
                                }
                            ],
                        )
                    except Exception as e:
                        self.logger.error(f"Error reindexing knowledge {knowledge.id}: {str(e)}")
                        continue

            return True
        except Exception as e:
            self.logger.error(f"Error reindexing knowledge: {str(e)}")
            raise

    def get_embeddings(self, knowledge_id: int, user_id: int) -> List[float]:
        """Get embeddings for a specific knowledge document"""
        try:
            # First verify ownership
            knowledge = self.get_knowledge(knowledge_id, user_id)
            if not knowledge:
                self.logger.warning(f"Knowledge {knowledge_id} not found or not accessible by user {user_id}")
                return []

            # Get embeddings from ChromaDB
            result = self.collection.get(
                where={"knowledge_id": str(knowledge_id)},
                include=["embeddings"]
            )

            # Handle numpy array properly
            if result and "embeddings" in result and len(result["embeddings"]) > 0:
                # Convert numpy array to list and get first embedding
                embeddings = result["embeddings"][0]
                return embeddings.tolist() if hasattr(embeddings, "tolist") else embeddings

            # If no embeddings found, generate new ones
            embedding = self._generate_embedding(knowledge.content)

            # Store in ChromaDB
            self.collection.upsert(
                ids=[str(knowledge_id)],
                embeddings=[embedding],
                documents=[knowledge.content],
                metadatas=[
                    {
                        "user_id": str(user_id),
                        "knowledge_id": str(knowledge_id),
                        "name": knowledge.name,
                        "identifier": knowledge.identifier,
                        "citation": knowledge.name
                    }
                ],
            )

            return embedding

        except Exception as e:
            self.logger.error(f"Error getting embeddings for knowledge {knowledge_id}: {str(e)}")
            return []

    def get_chunks_for_knowledge(self, knowledge_id):
        """
        Get all chunks for a specific knowledge document.
        Uses caching to avoid repeated database queries.
        
        Args:
            knowledge_id: ID of the knowledge document
            
        Returns:
            List of chunks with content and metadata
        """
        try:
            # Check cache first
            cache_key = f"chunks_{knowledge_id}"
            if cache_key in self._chunks_cache:
                self._cache_hits['chunks'] += 1
                return self._chunks_cache[cache_key]
                
            self._cache_misses['chunks'] += 1
            
            # Query ChromaDB for chunks with this knowledge_id
            results = self.collection.get(
                where={"knowledge_id": str(knowledge_id)},
                include=["documents", "metadatas", "embeddings"]
            )
            
            if not results or not results["documents"]:
                return []
                
            # Format chunks
            chunks = []
            for i, (doc, meta) in enumerate(zip(results["documents"], results["metadatas"])):
                chunk_id = f"{knowledge_id}_c{i}"
                chunks.append({
                    "id": chunk_id,
                    "content": doc,
                    "metadata": meta
                })
                
            # Cache the results
            self._chunks_cache[cache_key] = chunks
            
            # Manage cache size
            if len(self._chunks_cache) > self._chunks_cache_max_size:
                # Remove a random key
                self._chunks_cache.pop(next(iter(self._chunks_cache)))
                
            return chunks
        except Exception as e:
            self.logger.error(f"Error getting chunks for knowledge {knowledge_id}: {str(e)}")
            return []

    def get_citations_for_chunks(self, chunk_ids):
        """
        Get citation information for specific chunks.
        
        Args:
            chunk_ids: List of chunk IDs to retrieve citations for
            
        Returns:
            Dictionary mapping chunk IDs to citation information
        """
        try:
            if not chunk_ids:
                self.logger.warning("No chunk IDs provided for citation retrieval")
                return {}
                
            self.logger.info(f"Retrieving citations for {len(chunk_ids)} chunks")
            print(f"DEBUG: Retrieving citations for chunks: {chunk_ids}")
            
            # Extract knowledge IDs from chunk IDs (format: chunk_knowledge_id_index)
            knowledge_ids = set()
            chunk_to_knowledge_map = {}
            
            for chunk_id in chunk_ids:
                # Try to extract knowledge_id from chunk_id
                try:
                    if chunk_id.startswith("chunk_"):
                        parts = chunk_id.split("_")
                        if len(parts) >= 2:
                            knowledge_id = parts[1]
                            knowledge_ids.add(knowledge_id)
                            chunk_to_knowledge_map[chunk_id] = knowledge_id
                    else:
                        # For other formats, just use the chunk_id as is
                        chunk_to_knowledge_map[chunk_id] = chunk_id
                except Exception as e:
                    self.logger.warning(f"Could not parse knowledge_id from chunk_id {chunk_id}: {str(e)}")
                    continue
            
            print(f"DEBUG: Extracted knowledge IDs: {knowledge_ids}")
            
            # If we couldn't extract any knowledge IDs, return empty
            if not knowledge_ids:
                return {}
            
            # Query ChromaDB for chunks with these knowledge IDs
            citations = {}
            
            for knowledge_id in knowledge_ids:
                try:
                    # Query by knowledge_id
                    results = self.collection.get(
                        where={"knowledge_id": knowledge_id},
                        include=["metadatas"]
                    )
                    
                    if results and "metadatas" in results and results["metadatas"]:
                        # Use the first metadata entry for this knowledge_id
                        meta = results["metadatas"][0]
                        
                        # Extract citation information from metadata
                        citation_info = {
                            "source": meta.get("source", "Unknown Source"),
                            "citation": meta.get("citation", "Unknown Source"),
                            "knowledge_id": knowledge_id,
                        }
                        
                        # Add page/row information if available
                        if "page" in meta:
                            citation_info["page"] = meta["page"]
                        if "row" in meta:
                            citation_info["row"] = meta["row"]
                        
                        print(f"DEBUG: Found citation info for knowledge_id {knowledge_id}: {citation_info}")
                        
                        # Add this citation info to all chunks with this knowledge_id
                        for chunk_id, mapped_knowledge_id in chunk_to_knowledge_map.items():
                            if mapped_knowledge_id == knowledge_id:
                                citations[chunk_id] = citation_info
                                print(f"DEBUG: Added citation info for chunk {chunk_id}")
                except Exception as e:
                    self.logger.error(f"Error retrieving citation for knowledge_id {knowledge_id}: {str(e)}")
                    continue
            
            self.logger.info(f"Retrieved citation information for {len(citations)} chunks")
            print(f"DEBUG: Retrieved citation information for {len(citations)} chunks")
            return citations
        except Exception as e:
            self.logger.error(f"Error retrieving citations: {str(e)}")
            print(f"DEBUG: Error retrieving citations: {str(e)}")
            traceback.print_exc()
            return {}
            
    def format_citation(self, citation_info):
        """
        Format citation information into a human-readable string.
        
        Args:
            citation_info: Dictionary containing citation metadata
            
        Returns:
            Formatted citation string
        """
        try:
            if not citation_info:
                return "Unknown Source"
                
            # Start with the basic citation (usually the document name)
            citation = citation_info.get("citation", citation_info.get("source", "Unknown Source"))
            
            # Add page information if available
            if "page" in citation_info:
                citation = f"{citation}, Page {citation_info['page']}"
                
            # Add row information if available (for CSV files)
            if "row" in citation_info:
                citation = f"{citation}, Row {citation_info['row']}"
                
            return citation
        except Exception as e:
            self.logger.error(f"Error formatting citation: {str(e)}")
            return "Unknown Source"
            
    def get_citations_for_response(self, response_text, relevant_chunks):
        """
        Extract and format citations for a generated response based on relevant chunks.
        
        Args:
            response_text: The generated response text
            relevant_chunks: List of chunks used to generate the response
            
        Returns:
            Dictionary with citation information
        """
        try:
            self.logger.info(f"get_citations_for_response called with {len(relevant_chunks) if relevant_chunks else 0} chunks")
            
            if not relevant_chunks:
                self.logger.warning("No relevant chunks provided for citation generation")
                return {"citations": [], "has_citations": False}
            
            # Extract chunk IDs and metadata from relevant chunks
            chunk_ids = []
            chunk_metadata = {}
            
            for chunk in relevant_chunks:
                if isinstance(chunk, dict):
                    # Try to get the ID from the chunk
                    chunk_id = None
                    
                    if "id" in chunk:
                        chunk_id = chunk["id"]
                    elif "chunk_id" in chunk:
                        chunk_id = chunk["chunk_id"]
                    elif "metadata" in chunk and "chunk_id" in chunk["metadata"]:
                        chunk_id = chunk["metadata"]["chunk_id"]
                    else:
                        # Generate a temporary ID
                        chunk_id = f"temp_chunk_{len(chunk_ids)}"
                        
                    chunk_ids.append(chunk_id)
                    
                    # Store metadata directly if available
                    if "metadata" in chunk and isinstance(chunk["metadata"], dict):
                        chunk_metadata[chunk_id] = chunk["metadata"]
            
            # Get citation information for chunks
            citations_from_db = self.get_citations_for_chunks(chunk_ids)
            
            # Combine with direct metadata
            all_citations = {}
            for chunk_id in chunk_ids:
                if chunk_id in citations_from_db:
                    all_citations[chunk_id] = citations_from_db[chunk_id]
                elif chunk_id in chunk_metadata:
                    all_citations[chunk_id] = chunk_metadata[chunk_id]
            
            # Look for citation patterns in the response text
            # Use more specific patterns to avoid matching non-citation text
            citation_patterns = [
                r'\[(\d+)\]',  # [1], [2], etc.
                r'\[(Source|Reference|Citation|Document)\s+(\d+|[A-Za-z]+)\]',  # [Source 1], [Reference A]
                r'Source\s+(\d+|[A-Za-z]+)',  # Source 1 or Source A
                r'Reference\s+(\d+|[A-Za-z]+)',  # Reference 1
                r'Citation\s+(\d+|[A-Za-z]+)',  # Citation 1
                r'Document\s+(\d+|[A-Za-z]+)'   # Document 1
            ]
            
            # Clean the response text first - remove any sequences of [?] or [0]
            cleaned_response = re.sub(r'(\[\?+\]|\[0+\])+', '', response_text)
            
            found_citations = set()
            for pattern in citation_patterns:
                matches = re.finditer(pattern, cleaned_response)
                for match in matches:
                    # For patterns with multiple capture groups, use the last non-empty group
                    citation_text = None
                    if match.groups():
                        for group in match.groups():
                            if group:
                                citation_text = group
                    else:
                        citation_text = match.group(0)
                    
                    # Skip citations that are just question marks or zeros
                    if citation_text and not re.match(r'^[\?\[\]0]+$', citation_text):
                        found_citations.add(citation_text.strip())
            
            # Format the citations
            formatted_citations = []
            
            # First add citations found in the text
            for citation_text in found_citations:
                # Try to match citation text with our chunks
                matched = False
                for chunk_id, citation_info in all_citations.items():
                    citation_source = citation_info.get("citation", "")
                    if citation_text in citation_source or citation_source in citation_text:
                        formatted_citation = self.format_citation(citation_info)
                        formatted_citations.append({
                            "text": citation_text,
                            "source": formatted_citation,
                            "chunk_id": chunk_id,
                            "knowledge_id": citation_info.get("knowledge_id", ""),
                            "metadata": citation_info
                        })
                        matched = True
                        break
                
                # If no match, add as is
                if not matched:
                    formatted_citations.append({
                        "text": citation_text,
                        "source": citation_text,
                        "chunk_id": f"auto_citation_{len(formatted_citations)}",
                        "knowledge_id": "",
                        "metadata": {"citation": citation_text}
                    })
            
            # If no citations were found in text but we have chunks, add them as implicit citations
            if not formatted_citations and all_citations:
                for chunk_id, citation_info in all_citations.items():
                    formatted_citation = self.format_citation(citation_info)
                    formatted_citations.append({
                        "text": formatted_citation,
                        "source": formatted_citation,
                        "chunk_id": chunk_id,
                        "knowledge_id": citation_info.get("knowledge_id", ""),
                        "metadata": citation_info
                    })
            
            return {
                "citations": formatted_citations,
                "has_citations": len(formatted_citations) > 0
            }
        except Exception as e:
            self.logger.error(f"Error generating citations: {str(e)}")
            traceback.print_exc()
            return {"citations": [], "has_citations": False}

    def _process_file(self, knowledge, file):
        """
        Process a file and update the knowledge document with its content.
        Handles different file types and creates embeddings.
        
        Args:
            knowledge: The knowledge document to update
            file: The file to process
        """
        try:
            self.logger.info(f"Processing file {file.name} for knowledge {knowledge.id}")
            
            # Create a file-like object with the file content to avoid closed file issues
            file_info = {
                'name': file.name,
                'content_type': file.content_type,
                'size': file.size
            }
            
            # Read the file content into memory first
            try:
                if hasattr(file, 'read'):
                    file_content = file.read()
                    file_info['content'] = file_content
                else:
                    # If file is already closed, try to open it from the path
                    try:
                        with open(file.temporary_file_path(), 'rb') as f:
                            file_content = f.read()
                            file_info['content'] = file_content
                    except (AttributeError, FileNotFoundError):
                        # If we can't get the file content, raise an error
                        raise ValueError(f"Could not read file content for {file.name}")
            except Exception as e:
                self.logger.error(f"Error reading file content: {str(e)}")
                raise
            
            # Delegate to the content-based processing method
            return self._process_file_with_content(knowledge, file_info)
                
        except Exception as e:
            self.logger.error(f"Error processing file: {str(e)}")
            traceback.print_exc()
            # Update knowledge document with error status
            self.repository.update(knowledge.id, {
                "status": "error",
                "error_message": str(e),
            })

    def get_cache_stats(self):
        """Get statistics about cache performance"""
        stats = {
            'hits': self._cache_hits,
            'misses': self._cache_misses,
            'sizes': {
                'embedding': len(self._embedding_cache),
                'search': len(self._search_results_cache),
                'chunks': len(self._chunks_cache)
            }
        }
        
        # Calculate hit rates
        hit_rates = {}
        for cache_type in self._cache_hits.keys():
            total = self._cache_hits[cache_type] + self._cache_misses[cache_type]
            hit_rates[cache_type] = self._cache_hits[cache_type] / total if total > 0 else 0
            
        stats['hit_rates'] = hit_rates
        return stats

    def _process_pptx(self, file_path, document_name):
        """
        Process a PowerPoint presentation (PPTX/PPT) and extract its content.
        
        Args:
            file_path: Path to the PowerPoint presentation
            document_name: Name of the document
            
        Returns:
            Tuple of (full_content, chunks, metadata)
        """
        try:
            # Import pptx library here to avoid dependency issues
            from pptx import Presentation
            
            full_content = ""
            chunks = []
            metadata = {"source": document_name, "type": "pptx"}
            
            # Open the presentation
            presentation = Presentation(file_path)
            
            # Extract text from slides
            slides_content = []
            
            for slide_index, slide in enumerate(presentation.slides):
                slide_text = []
                slide_number = slide_index + 1
                
                # Get slide title if available
                if slide.shapes.title and slide.shapes.title.text:
                    slide_text.append(f"Slide {slide_number} - Title: {slide.shapes.title.text}")
                else:
                    slide_text.append(f"Slide {slide_number}")
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip if it's the title we already added
                        if shape == slide.shapes.title:
                            continue
                        slide_text.append(shape.text)
                
                # Extract text from tables in the slide
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                            if row_text.strip():
                                slide_text.append(row_text)
                
                # Combine all text from this slide
                slide_content = "\n".join(slide_text)
                slides_content.append(slide_content)
                
                # Create a chunk for each slide with citation metadata
                chunk_id = f"{document_name}_slide{slide_number}"
                chunks.append({
                    "id": chunk_id,
                    "content": slide_content,
                    "metadata": {
                        "source": document_name,
                        "slide": slide_number,
                        "citation": f"{document_name}, Slide {slide_number}"
                    }
                })
            
            # Combine all slides
            full_content = "\n\n".join(slides_content)
            
            return full_content, chunks, metadata
        except ImportError:
            self.logger.error("python-pptx is not installed. Please install it to process PPTX files.")
            raise Exception("PPTX processing library not available")
        except Exception as e:
            self.logger.error(f"Error processing PPTX: {str(e)}")
            raise
