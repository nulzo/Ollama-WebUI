import os
import sys
import django
from django.conf import settings
import tempfile
import base64
import json
import time

# Set up Django environment
os.environ.setdefault("DJANGO_SETTINGS_MODULE", "settings.settings")
django.setup()

from features.knowledge.services.knowledge_service import KnowledgeService
from features.completions.services.chat_service import ChatService
from features.authentication.models import CustomUser

def create_test_files():
    """Create test files of different types for testing"""
    test_files = {}
    
    # Create a text file
    with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as f:
        f.write(b"This is a test text file.\nIt contains information about RAG testing.\nRetrieval Augmented Generation is a technique to enhance LLM responses with external knowledge.")
        test_files['txt'] = {'path': f.name, 'name': 'test.txt', 'content_type': 'text/plain'}
    
    # Create a CSV file
    with tempfile.NamedTemporaryFile(suffix='.csv', delete=False) as f:
        f.write(b"Name,Description,Value\nRAG,Retrieval Augmented Generation,High\nLLM,Large Language Model,Medium\nAPI,Application Programming Interface,Low")
        test_files['csv'] = {'path': f.name, 'name': 'test.csv', 'content_type': 'text/csv'}
    
    # Create a Markdown file
    with tempfile.NamedTemporaryFile(suffix='.md', delete=False) as f:
        f.write(b"# RAG Testing\n\n## Introduction\nThis is a test markdown file for RAG functionality.\n\n## Features\n- Retrieval\n- Augmentation\n- Generation")
        test_files['md'] = {'path': f.name, 'name': 'test.md', 'content_type': 'text/markdown'}
    
    # Create an HTML file
    with tempfile.NamedTemporaryFile(suffix='.html', delete=False) as f:
        f.write(b"""<!DOCTYPE html>
<html>
<head>
    <title>RAG Test Document</title>
</head>
<body>
    <h1>RAG Testing</h1>
    <p>This is a test HTML document for testing RAG functionality.</p>
    <h2>What is RAG?</h2>
    <p>Retrieval Augmented Generation (RAG) is a technique that enhances LLM responses with external knowledge.</p>
</body>
</html>""")
        test_files['html'] = {'path': f.name, 'name': 'test.html', 'content_type': 'text/html'}
    
    # Create a JSON file
    with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as f:
        json_data = {
            "title": "RAG Testing",
            "description": "Test JSON file for RAG functionality",
            "concepts": [
                {"name": "Retrieval", "description": "Finding relevant information"},
                {"name": "Augmentation", "description": "Adding context to prompts"},
                {"name": "Generation", "description": "Creating responses with LLMs"}
            ]
        }
        f.write(json.dumps(json_data).encode('utf-8'))
        test_files['json'] = {'path': f.name, 'name': 'test.json', 'content_type': 'application/json'}
    
    # Create a PowerPoint file
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        pptx_path = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False).name
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "RAG Testing with PowerPoint"
        subtitle.text = "Testing PowerPoint support in RAG functionality"
        
        # Add a content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = "What is RAG?"
        
        # Add bullet points
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Retrieval Augmented Generation"
        
        p = tf.add_paragraph()
        p.text = "Enhances LLM responses with external knowledge"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Improves accuracy and reduces hallucinations"
        p.level = 1
        
        # Add a table slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "RAG Components"
        
        # Add a table
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(1.5)
        shape = slide.shapes.add_table(3, 2, x, y, cx, cy)
        table = shape.table
        
        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(3)
        
        # Add header row
        cell = table.cell(0, 0)
        cell.text = "Component"
        cell = table.cell(0, 1)
        cell.text = "Description"
        
        # Add data rows
        cell = table.cell(1, 0)
        cell.text = "Retriever"
        cell = table.cell(1, 1)
        cell.text = "Finds relevant documents"
        
        cell = table.cell(2, 0)
        cell.text = "Generator"
        cell = table.cell(2, 1)
        cell.text = "Creates responses using context"
        
        # Save the presentation
        prs.save(pptx_path)
        test_files['pptx'] = {'path': pptx_path, 'name': 'test.pptx', 'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'}
        print("Created PowerPoint test file")
    except Exception as e:
        print(f"Error creating PowerPoint test file: {str(e)}")
    
    return test_files

def test_rag_functionality():
    """Test the RAG functionality with different file types"""
    print("Testing RAG Functionality...")
    
    # Initialize services
    knowledge_service = KnowledgeService()
    chat_service = ChatService()
    
    # Get a test user
    try:
        user = CustomUser.objects.first()
        if not user:
            print("No users found in the database. Please create a user first.")
            return
        print(f"Using test user: {user.username} (ID: {user.id})")
    except Exception as e:
        print(f"Error getting test user: {str(e)}")
        return
    
    # Create test files
    test_files = create_test_files()
    print(f"Created {len(test_files)} test files")
    
    # Process each test file
    knowledge_ids = []
    for file_type, file_info in test_files.items():
        try:
            print(f"\nProcessing {file_type.upper()} file: {file_info['name']}")
            
            # Read file content
            with open(file_info['path'], 'rb') as f:
                file_content = f.read()
            
            # Create file info dictionary
            file_data = {
                'name': file_info['name'],
                'content_type': file_info['content_type'],
                'size': len(file_content),
                'content': file_content
            }
            
            # Create knowledge document
            knowledge = knowledge_service.create_knowledge({
                "name": f"Test {file_type.upper()} Knowledge",
                "identifier": f"test_{file_type}_knowledge_{django.utils.timezone.now().timestamp()}",
                "content": "",  # Will be populated by the document processor
                "file_path": file_info['name'],
                "file_size": len(file_content),
                "file_type": file_info['content_type'],
                "status": "processing",
            }, user)
            
            print(f"Created knowledge document: {knowledge.name} (ID: {knowledge.id})")
            
            # Process the file
            knowledge_service._process_file_with_content(knowledge, file_data)
            
            # Verify the knowledge document was updated
            knowledge = knowledge_service.get_knowledge(knowledge.id, user.id)
            print(f"Knowledge status: {knowledge.status}")
            print(f"Content length: {len(knowledge.content)}")
            
            # Get chunks
            chunks = knowledge_service.get_chunks_for_knowledge(knowledge.id)
            print(f"Found {len(chunks)} chunks")
            
            # Store knowledge ID for later testing
            knowledge_ids.append(knowledge.id)
            
        except Exception as e:
            print(f"Error processing {file_type} file: {str(e)}")
            import traceback
            traceback.print_exc()
    
    # Test retrieval with a query
    if knowledge_ids:
        try:
            print("\nTesting retrieval with a query...")
            query = "What is RAG?"
            
            # Test semantic search
            print("\nTesting semantic search...")
            results = knowledge_service._semantic_search(query, user.id, max_results=3)
            print(f"Found {len(results)} relevant documents with semantic search")
            for i, result in enumerate(results):
                print(f"Result {i+1}:")
                print(f"  Content: {result['content'][:100]}...")
                print(f"  Citation: {result['metadata'].get('citation', 'Unknown')}")
                print(f"  Similarity: {result['similarity']:.4f}")
            
            # Test keyword search
            print("\nTesting keyword search...")
            results = knowledge_service._keyword_search(query, user.id, max_results=3)
            print(f"Found {len(results)} relevant documents with keyword search")
            for i, result in enumerate(results):
                print(f"Result {i+1}:")
                print(f"  Content: {result['content'][:100]}...")
                print(f"  Citation: {result['metadata'].get('citation', 'Unknown')}")
                print(f"  Similarity: {result['similarity']:.4f}")
            
            # Test hybrid search
            print("\nTesting hybrid search...")
            results = knowledge_service._hybrid_search(query, user.id, max_results=3)
            print(f"Found {len(results)} relevant documents with hybrid search")
            for i, result in enumerate(results):
                print(f"Result {i+1}:")
                print(f"  Content: {result['content'][:100]}...")
                print(f"  Citation: {result['metadata'].get('citation', 'Unknown')}")
                print(f"  Similarity: {result['similarity']:.4f}")
                print(f"  Search type: {result['search_type']}")
                print(f"  Combined score: {result.get('combined_score', 0):.4f}")
            
            # Test context preparation in chat service
            print("\nTesting context preparation...")
            context = chat_service._prepare_context(query, user.id, knowledge_ids=knowledge_ids)
            print(f"Context length: {len(context)}")
            print(f"Context preview: {context[:200]}...")
            
            # Test caching
            print("\nTesting caching...")
            # First call should be a cache miss
            start_time = time.time()
            knowledge_service._hybrid_search(query, user.id, max_results=3)
            first_call_time = time.time() - start_time
            
            # Second call should be a cache hit
            start_time = time.time()
            knowledge_service._hybrid_search(query, user.id, max_results=3)
            second_call_time = time.time() - start_time
            
            print(f"First call time: {first_call_time:.4f}s")
            print(f"Second call time: {second_call_time:.4f}s")
            print(f"Speedup: {first_call_time / second_call_time if second_call_time > 0 else 'N/A'}x")
            
            # Get cache stats
            cache_stats = knowledge_service.get_cache_stats()
            print("\nCache statistics:")
            print(f"  Hits: {cache_stats['hits']}")
            print(f"  Misses: {cache_stats['misses']}")
            print(f"  Hit rates: {cache_stats['hit_rates']}")
            print(f"  Cache sizes: {cache_stats['sizes']}")
            
            # Clean up - delete the test knowledge documents
            print("\nCleaning up...")
            for knowledge_id in knowledge_ids:
                try:
                    knowledge_service.delete_knowledge(knowledge_id, user.id)
                    print(f"Deleted knowledge document: {knowledge_id}")
                except Exception as e:
                    print(f"Error deleting knowledge {knowledge_id}: {str(e)}")
            
            # Clean up test files
            for file_info in test_files.values():
                try:
                    os.unlink(file_info['path'])
                except Exception as e:
                    print(f"Error deleting file {file_info['path']}: {str(e)}")
            
        except Exception as e:
            print(f"Error testing retrieval: {str(e)}")
            import traceback
            traceback.print_exc()

if __name__ == "__main__":
    test_rag_functionality() 